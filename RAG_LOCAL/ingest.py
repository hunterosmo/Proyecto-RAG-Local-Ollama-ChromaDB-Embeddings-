# ingest.py (versión con PDF + Office + txt/md, optimizada y segura)
import uuid
import gc
from pathlib import Path
from datetime import datetime

import chromadb
from sentence_transformers import SentenceTransformer

from config import DOCS_DIR, CHROMA_DIR, EMBEDDING_MODEL_NAME, CHUNK_SIZE, CHUNK_OVERLAP

# Librerías para formatos específicos
from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
from openpyxl import load_workbook

EXTENSIONES_SOPORTADAS = [".txt", ".md", ".pdf", ".docx", ".pptx", ".xlsx"]

# 🔢 Tamaño de lote para embeddings (ajusta si quieres usar menos RAM aún)
BATCH_SIZE = 16


def chunk_text(text: str, max_chars: int, overlap: int):
    """
    Divide un texto largo en chunks con solapamiento por caracteres,
    usando un cálculo seguro basado en pasos fijos.
    No hay bucles infinitos.
    """
    if max_chars <= 0:
        raise ValueError("max_chars debe ser > 0")

    if overlap < 0:
        overlap = 0
    if overlap >= max_chars:
        print(f"⚠ overlap ({overlap}) >= max_chars ({max_chars}), ajustando overlap automáticamente...")
        overlap = max_chars - 1

    length = len(text)
    step = max_chars - overlap  # cuánto avanzamos cada vez
    if step <= 0:
        step = 1

    chunks = []

    print(f"   · Texto total: {length} caracteres")
    print(f"   · max_chars={max_chars}, overlap={overlap}, step={step}")

    for start in range(0, length, step):
        end = min(start + max_chars, length)
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)

    print(f"   · Chunks generados: {len(chunks)}")
    return chunks


def load_txt_md(path: Path) -> str:
    with path.open("r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def load_pdf(path: Path) -> str:
    """Extrae texto de un PDF página por página."""
    texto = []
    reader = PdfReader(str(path))
    for i, page in enumerate(reader.pages):
        try:
            page_text = page.extract_text() or ""
        except Exception:
            page_text = ""
        if page_text.strip():
            texto.append(f"[Página {i+1}]\n{page_text}")
    return "\n\n".join(texto)


def load_docx(path: Path) -> str:
    """Extrae texto de un documento .docx (Word)."""
    doc = DocxDocument(str(path))
    parts = []
    for para in doc.paragraphs:
        txt = para.text.strip()
        if txt:
            parts.append(txt)
    return "\n".join(parts)


def load_pptx(path: Path) -> str:
    """Extrae texto de un .pptx (PowerPoint)."""
    prs = Presentation(str(path))
    parts = []
    for i, slide in enumerate(prs.slides):
        slide_parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                txt = shape.text.strip()
                if txt:
                    slide_parts.append(txt)
        if slide_parts:
            parts.append(f"[Diapositiva {i+1}]\n" + "\n".join(slide_parts))
    return "\n\n".join(parts)


def load_xlsx(path: Path) -> str:
    """Extrae texto de un .xlsx (Excel) recorriendo hojas y celdas."""
    wb = load_workbook(filename=str(path), data_only=True)
    parts = []
    for sheet in wb.worksheets:
        sheet_parts = []
        for row in sheet.iter_rows(values_only=True):
            vals = [str(c) for c in row if c is not None]
            if vals:
                sheet_parts.append(" | ".join(vals))
        if sheet_parts:
            parts.append(f"[Hoja: {sheet.title}]\n" + "\n".join(sheet_parts))
    return "\n\n".join(parts)


def load_file(path: Path) -> str:
    """Detecta el tipo de archivo y llama al loader correspondiente."""
    ext = path.suffix.lower()

    if ext in [".txt", ".md"]:
        return load_txt_md(path)
    if ext == ".pdf":
        return load_pdf(path)
    if ext == ".docx":
        return load_docx(path)
    if ext == ".pptx":
        return load_pptx(path)
    if ext == ".xlsx":
        return load_xlsx(path)

    return ""


def main():
    DOCS_DIR.mkdir(parents=True, exist_ok=True)

    print(f"📂 Carpeta de documentos: {DOCS_DIR}")
    print(f"🔍 Buscando archivos con estas extensiones: {', '.join(EXTENSIONES_SOPORTADAS)}")

    files = [p for p in DOCS_DIR.glob("**/*") if p.suffix.lower() in EXTENSIONES_SOPORTADAS]

    if not files:
        print("⚠ No se encontraron archivos compatibles en la carpeta docs.")
        print("   Coloca tus PDF, DOCX, PPTX, XLSX, TXT o MD en D:\\RAG_LOCAL\\docs y vuelve a ejecutar este script.")
        return

    print(f"✅ Encontrados {len(files)} archivo(s).")

    print(f"🧠 Cargando modelo de embeddings: {EMBEDDING_MODEL_NAME}...")
    embedder = SentenceTransformer(EMBEDDING_MODEL_NAME)

    print(f"📚 Iniciando Chroma en: {CHROMA_DIR}")
    client = chromadb.PersistentClient(path=str(CHROMA_DIR))
    collection = client.get_or_create_collection(name="docs")

    for file_path in files:
        print(f"\n📄 Procesando: {file_path.name} ({file_path.suffix.lower()})")

        # Protección por tamaño (ej. > 200 MB)
        file_size_mb = file_path.stat().st_size / (1024 * 1024)
        print(f"   · Tamaño archivo: {file_size_mb:.2f} MB")
        if file_size_mb > 200:
            print("   ⚠ Archivo demasiado grande, se omite por seguridad.")
            continue

        try:
            text = load_file(file_path)
        except Exception as e:
            print(f"   ⚠ Error leyendo el archivo: {e}")
            continue

        if not text or not text.strip():
            print("   (Archivo sin texto útil, se omite)")
            continue

        chunks = chunk_text(text, max_chars=CHUNK_SIZE, overlap=CHUNK_OVERLAP)
        if not chunks:
            print("   (No se generaron chunks, se omite)")
            continue

        # Metadatos base
        rel_path = file_path.relative_to(DOCS_DIR)
        folder = str(rel_path.parent) if rel_path.parent != Path('.') else ""
        ext = file_path.suffix.lower()
        mtime = file_path.stat().st_mtime
        mdate = datetime.fromtimestamp(mtime).strftime("%Y-%m-%d")

        chunk_index_offset = 0

        # Procesar en lotes
        for start in range(0, len(chunks), BATCH_SIZE):
            batch_chunks = chunks[start:start + BATCH_SIZE]

            batch_ids = []
            batch_metadatas = []
            for i, _ in enumerate(batch_chunks):
                idx = chunk_index_offset + i
                batch_ids.append(str(uuid.uuid4()))
                batch_metadatas.append({
                    "source": str(file_path),
                    "chunk_index": idx,
                    "ext": ext,
                    "folder": folder,
                    "date": mdate,
                })

            chunk_index_offset += len(batch_chunks)

            print(f"   · Lote de {len(batch_chunks)} chunks → generando embeddings...")
            embeddings = embedder.encode(batch_chunks, show_progress_bar=False).tolist()

            collection.add(
                documents=batch_chunks,
                embeddings=embeddings,
                metadatas=batch_metadatas,
                ids=batch_ids
            )
            print(f"   · Lote guardado en Chroma.")

            # Liberar memoria del lote
            del batch_chunks, batch_ids, batch_metadatas, embeddings
            gc.collect()

    print("\n✅ Ingesta completada. Tu base vectorial está lista.")


if __name__ == "__main__":
    main()
